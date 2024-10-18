# base imports for node
import logging
import knime.extension as knext

# imports for translation logic
import os
import time
import requests
import re
from pptx import Presentation
from pptx.util import Pt
import urllib.parse

LOGGER = logging.getLogger(__name__)


# 3 nodes
# - use existing prompt node
# - basic parser for powerpoint (just output text)
# - translation logic (new powerpoint file) (node for injecting text)
# custom port ?

# jenkins -> runs node to test

# Global variables for tracking use for rate limiting in translation
START_TIME = time.time()
REQUEST_COUNT = 0
TOKEN_COUNT = 0

@knext.node(name="Powerpoint Translator Test", node_type=knext.NodeType.MANIPULATOR, icon_path="icon.png", category="/")
class TemplateNode:
    """Powerpoint file translator using OpenAI.
    This node translates files passed using OpenAI's API. You will need a key to pass to this node using the credentials configuration node.
    """

    input_file = knext.LocalPathParameter(
        label="Input PowerPoint File",
        description="Select the PowerPoint file to translate.",
    )

    output_file_name = knext.StringParameter(
        label="Output File Name",
        description="Choose the base file name. After translation, it will have the target larguage appended to it. (myfile.pptx -> German_myfile.pptx)",
        default_value=""
    )

    target_language = knext.StringParameter(
        label="Target Language",
        description="Choose the language for translation.",
        default_value="English"
    )

    credential_param = knext.StringParameter(
        label="Open AI key",
        description="Choose one of the connected credentials (Pass key through credential config password field)",
        choices=lambda a: knext.DialogCreationContext.get_credential_names(a)
    )

    max_requests = knext.IntParameter(
        label="Max Requests Per Minute",
        description="Maximum number of API requests per minute.",
        is_advanced=True,
        default_value=3500
    )

    max_tokens = knext.IntParameter(
        label="Max Tokens Per Minute",
        description="Maximum tokens allowed per minute.",
        is_advanced=True,
        default_value=200000
    )

    # might take this out
    output_file = knext.LocalPathParameter(
        label="Output Directory (Optional)",
        description="Optional path for saving the translated file.",
        is_advanced=True,
        placeholder_text="Leave empty to use a temporary directory."
    )

    def configure(self, ctx: knext.ConfigurationContext):
        if not ctx.get_credential_names():
            raise knext.InvalidParametersError("No credentials provided.")
        if not self.credential_param:
            raise knext.InvalidParametersError("Credentials not selected.")
        return None

    def execute(self, ctx: knext.ExecutionContext):
        # Fetch credentials from KNIME node
        credentials = ctx.get_credentials(self.credential_param)
        self.api_key = credentials.password 

        self.input_file_path = ctx.flow_variables.get("input_file", self.input_file)
        self.target_language = ctx.flow_variables.get("target_language", self.target_language)
        self.output_file = ctx.flow_variables.get("output_file", self.output_file)
        self.max_tokens = ctx.flow_variables.get("max_tokens", self.max_tokens)
        self.max_requests = ctx.flow_variables.get("max_requests", self.max_requests)
        self.output_name = ctx.flow_variables.get("output_file_name",self.output_file_name)

        self.input_file_path = self.resolve_file_path(self.input_file_path)

        translated_file_path = self.process_presentation()

        if not translated_file_path:
            LOGGER.warning("Translation failed. Output file path might not be valid.")

        output_path = self.create_output_file(os.path.basename(translated_file_path))
        ctx.flow_variables["translated_pptx_path"] = output_path

    def create_output_file(self, filename):
        """
        Create the output file path based on user selection or fallback to the Desktop.
        Use the base name of the input file, and append the language for the translated file.
        """
        output_filename = f"{self.target_language}_translated_{self.output_name}.pptx"
        
        if self.output_file:  # Use the user-specified output directory if provided
            output_path = os.path.join(self.output_file, output_filename)
        else:
            # Hardcoded path to the desktop for now
            output_path = os.path.join("/Users/thorlandstrom/Desktop", output_filename)
        
        return output_path


    def resolve_file_path(self, uri):
        decoded_uri = urllib.parse.unquote(uri)
        if decoded_uri.startswith("knime://knime.workflow"):
            relative_path = decoded_uri.replace("knime://knime.workflow", "")
            return os.path.join("", relative_path.lstrip("/"))
        return decoded_uri

    def check_rate_limit(self, tokens):
        """
        Ensure that API rate limits are not exceeded.
        """
        global REQUEST_COUNT, TOKEN_COUNT, START_TIME

        REQUEST_COUNT += 1
        TOKEN_COUNT += tokens

        elapsed_time = time.time() - START_TIME
        if elapsed_time >= 60:
            REQUEST_COUNT, TOKEN_COUNT, START_TIME = 0, 0, time.time()

        if REQUEST_COUNT >= self.max_requests or TOKEN_COUNT >= self.max_tokens:
            wait_time = 60 - elapsed_time
            print(f"Rate limit reached. Waiting.")
            time.sleep(wait_time)
            REQUEST_COUNT, TOKEN_COUNT, START_TIME = 0, 0, time.time()

    def translate_text(self, text):
        """
        Use OpenAI API to translate a given text. Logs the API request and response.
        """
        if not text.strip() or not re.search(r'[a-zA-Z0-9]', text):
            LOGGER.info(f"Skipping non-translatable text: {text}")
            return text

        tokens_needed = len(text.split())
        self.check_rate_limit(tokens_needed)
        
        url = "https://api.openai.com/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }

        # prompt done with some other node 

        system_message = f"""You are a translator specializing in PowerPoint presentations. 
        Your task is to translate text to {self.target_language}. 
        For image attributions or license information, keep proper nouns, abbreviations, and license codes unchanged. 
        Translate only the surrounding text.
        If the text looks like it should not be translated, then leave it as is (such as dates, math, equations, etc.).
        IMPORTANT: Your response must be in the following format:
        [START_TRANSLATION]
        Your translated text here
        [END_TRANSLATION]
        Any explanations or notes should be outside these tags."""
        user_message = f"Translate the following text to {self.target_language}:\n\n{text}"

        body = {
            "model": "gpt-4o-mini", # change model here if it becomes deprecated
            "messages": [
                {"role": "system", "content": system_message},
                {"role": "user", "content": user_message}
            ],
            "max_tokens": 1000,
            "n": 1,
            "temperature": 0.1
        }

        try:
            response = requests.post(url, headers=headers, json=body)
            response.raise_for_status()  # Raise an error for non-2xx status codes
        except requests.exceptions.RequestException as e:
            LOGGER.error(f"Error during translation API call: {e}")
            return text

        if response.status_code == 200:
            content = response.json()['choices'][0]['message']['content'].strip()
            LOGGER.debug(f"Translated content: {content}")
            
            start_tag, end_tag = "[START_TRANSLATION]", "[END_TRANSLATION]"
            start_index, end_index = content.find(start_tag) + len(start_tag), content.find(end_tag)
            
            if start_index != -1 and end_index != -1:
                return content[start_index:end_index].strip()
            return content  # Return the full content if tags not found
        else:
            LOGGER.warning(f"Translation failed with status code: {response.status_code}")
            return text

    def adjust_font_size(self, run, original_text, translated_text):
        """
        Adjust font size to prevent overflow.
        """
        original_length, translated_length = len(original_text), len(translated_text)
        if run.font.size:
            current_size = run.font.size.pt
            scale_factor = min(original_length / translated_length, 1) if translated_length > original_length else 1
            new_size = max(10, min(current_size * scale_factor, 400))
            run.font.size = Pt(new_size)

    def translate_shape_text(self, shape):
        """
        Translate the text in a PowerPoint shape.
        """
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() and re.search(r'[a-zA-Z0-9]', run.text):
                        original_text = run.text
                        translated_text = self.translate_text(original_text)
                        if translated_text != original_text:
                            self.adjust_font_size(run, original_text, translated_text)
                            run.text = translated_text

    def translate_table(self, table):
        """
        Translate text inside PowerPoint tables.
        """
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip() and re.search(r'[a-zA-Z0-9]', run.text):
                            original_text = run.text
                            translated_text = self.translate_text(original_text)
                            if translated_text != original_text:
                                self.adjust_font_size(run, original_text, translated_text)
                                run.text = translated_text

    def process_shapes(self, shapes):
        """
        Recursively process and translate all shapes in the presentation.
        """
        for shape in shapes:
            if shape.has_text_frame:
                self.translate_shape_text(shape)
            elif shape.has_table:
                self.translate_table(shape.table)
            elif hasattr(shape, 'shapes'):  # Grouped shapes
                self.process_shapes(shape.shapes)

    def process_presentation(self):
        """
        Process the entire PowerPoint presentation and translate text.
        """
        print(f"Opening {self.input_file_path}")

        try:
            presentation = Presentation(self.input_file_path)
        except Exception as e:
            print(f"Error opening file: {e}")
            return None

        # tqdm not needed anymore
        '''
        slide_count = len(presentation.slides)
        with tqdm(total=slide_count, desc="Translating", unit="slide") as pbar:
            for slide in presentation.slides:
                self.process_shapes(slide.shapes)
                pbar.update(1)
        '''

        for slide in presentation.slides:
                self.process_shapes(slide.shapes)
    
        filename = os.path.basename(self.input_file_path)
        output_path = self.create_output_file(filename)

        try:
            presentation.save(output_path)  # Save the translated file using the full path
            print(f"Translated file saved as {output_path}")
            return output_path
        except Exception as e:
            print(f"Error saving file: {e}")
            return None
